from typing import List, Dict
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from dataclasses import dataclass

from template_config import (
    ORANGE, WHITE, BLACK,
    TITLE_FONT, BODY_FONT,
    TITLE_FONT_SIZE, SUBTITLE_FONT_SIZE,
    JEP_TITLE_FONT_SIZE,
    apply_title_slide_layout, apply_jep_slide_layout, apply_example_slide_layout
)

@dataclass
class JEP:
    number: str
    title: str
    description: str = ""
    examples: List[Dict[str, str]] = None
    
    def __post_init__(self):
        # Ensure examples is always a list
        if self.examples is None:
            self.examples = []

class JDKRelease:
    def __init__(self, version: str, release_date: str, tagline: str):
        self.version = version
        self.release_date = release_date
        self.tagline = tagline
        self.jeps = []
    
    def add_jep(self, jep: JEP):
        """Add a JEP to the release"""
        self.jeps.append(jep)

class JDKPresentationGenerator:
    def __init__(self):
        self.prs = Presentation()
        self.prs.slide_width = 9144000  # 10 inches in EMU (16:9 aspect ratio)
        self.prs.slide_height = 5143500  # 5.625 inches in EMU
        self.current_slide = None
    
    def create_title_slide(self, release: JDKRelease):
        """Create the title slide for the JDK release"""
        slide_layout = self.prs.slide_layouts[0]  # Title Slide layout
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        
        title = f"JAVA {release.version}"
        subtitle = f"{release.tagline} (Release date {release.release_date})"
        
        # Pass the slide height from the presentation
        slide_height = self.prs.slide_height
        apply_title_slide_layout(self.current_slide, title, subtitle, slide_height=slide_height)
    
    def create_jep_slide(self, jep: JEP):
        """Create a slide for a JEP"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        
        # Add JEP number and title with centered layout
        slide_height = self.prs.slide_height
        apply_jep_slide_layout(
            self.current_slide,
            jep.number,
            jep.title,
            slide_height=slide_height
        )
        
        # Add example slides if they exist
        if jep.examples:
            for i, example in enumerate(jep.examples, 1):
                # Combine title and content for the example slide
                example_title = f"JEP {jep.number} Example {i}"
                example_content = f"{example.get('title', '')}\n\n{example.get('content', '')}"
                self.create_example_slide(example_title, example_content)
    
    def create_example_slide(self, title: str, content: str):
        """Create a slide with an example"""
        slide_layout = self.prs.slide_layouts[1]  # Title and Content layout
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        slide_height = self.prs.slide_height
        apply_example_slide_layout(self.current_slide, title, content, slide_height=slide_height)
    
    def generate_jdk_presentation(self, release: JDKRelease, output_path: str):
        """Generate the complete JDK release presentation"""
        # Create title slide
        self.create_title_slide(release)
        
        # Create slides for each JEP
        for jep in release.jeps:
            self.create_jep_slide(jep)
        
        # Save the presentation
        self.save_presentation(output_path)
    
    def save_presentation(self, filename: str):
        """Save the presentation to a file."""
        self.prs.save(filename)

    def add_title_slide(self, title, subtitle=''):
        """Add a title slide to the presentation."""
        slide_layout = self.prs.slide_layouts[0]
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        title_shape = self.current_slide.shapes.title
        subtitle_shape = self.current_slide.placeholders[1]
        
        title_shape.text = title
        subtitle_shape.text = subtitle

    def add_content_slide(self, title, content):
        """Add a content slide with title and bullet points."""
        slide_layout = self.prs.slide_layouts[1]
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        title_shape = self.current_slide.shapes.title
        body_shape = self.current_slide.placeholders[1]
        
        title_shape.text = title
        tf = body_shape.text_frame
        tf.text = content

    def add_table_slide(self, title, data):
        """Add a slide with a table."""
        slide_layout = self.prs.slide_layouts[5]
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        title_shape = self.current_slide.shapes.title
        
        title_shape.text = title
        
        # Create table
        rows = len(data) + 1
        cols = len(data.columns)
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(1.5)
        
        table = self.current_slide.shapes.add_table(rows, cols, left, top, width, height).table
        
        # Add headers
        for idx, col in enumerate(data.columns):
            table.cell(0, idx).text = col
        
        # Add data
        for idx, row in enumerate(data.itertuples(), 1):
            for col_idx, value in enumerate(row[1:], 0):
                table.cell(idx, col_idx).text = str(value)

    def add_image_slide(self, title, image_url):
        """Add a slide with an image."""
        slide_layout = self.prs.slide_layouts[6]
        self.current_slide = self.prs.slides.add_slide(slide_layout)
        title_shape = self.current_slide.shapes.title
        
        title_shape.text = title
        
        # Download image
        response = requests.get(image_url)
        with open('temp_image.jpg', 'wb') as f:
            f.write(response.content)
        
        # Add image
        left = Inches(1)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)
        self.current_slide.shapes.add_picture('temp_image.jpg', left, top, width, height)
        os.remove('temp_image.jpg')

    def save_presentation(self, filename):
        """Save the presentation to a file."""
        self.prs.save(filename)

def generate_presentation_from_web_data(url):
    """Generate a presentation from web data."""
    # Fetch and parse web data
    response = requests.get(url)
    soup = BeautifulSoup(response.text, 'html.parser')
    
    # Create presentation
    generator = PresentationGenerator()
    
    # Add title slide
    generator.add_title_slide("Web Data Presentation", "Generated from online content")
    
    # Add content slides based on web data
    # This is a simplified example - you would need to customize based on actual web content
    generator.add_content_slide("Source Content", "Data extracted from: " + url)
    
    # Add any tables or images if present
    tables = soup.find_all('table')
    if tables:
        df = pd.read_html(str(tables[0]))[0]
        generator.add_table_slide("Extracted Table", df)
    
    images = soup.find_all('img')
    if images:
        img_url = images[0]['src']
        generator.add_image_slide("Extracted Image", img_url)
    
    # Save presentation
    generator.save_presentation('web_data_presentation.pptx')
    print("Presentation generated successfully!")

if __name__ == "__main__":
    # Example usage
    url = "https://example.com"  # Replace with actual URL
    generate_presentation_from_web_data(url)
