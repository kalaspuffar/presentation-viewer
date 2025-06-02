from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor

# Colors
ORANGE = RGBColor(255, 87, 34)  # FF5722
WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)

# Fonts
TITLE_FONT = 'Alfa Slab One'
BODY_FONT = 'Roboto'
TITLE_FONT_SIZE = Pt(48)
SUBTITLE_FONT_SIZE = Pt(24)
JEP_TITLE_FONT_SIZE = Pt(24)

# Layout
TITLE_TOP_MARGIN = Inches(2)
CONTENT_LEFT_MARGIN = Inches(1)
CONTENT_TOP_MARGIN = Inches(2)
CONTENT_WIDTH = Inches(8.5)

# Slide dimensions (16:9 aspect ratio)
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

def apply_title_slide_layout(slide, title_text, subtitle_text, slide_height=Inches(5.625)):
    """Apply layout for title slide"""
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ORANGE
    
    # Clear any existing placeholders first
    for shape in slide.shapes:
        if not (shape.has_text_frame and shape.text.strip()):
            sp = shape._element
            sp.getparent().remove(sp)
    
    # Calculate total height needed for title and subtitle
    title_height = Inches(1.5)  # Approximate height for title
    subtitle_height = Inches(0.7)  # Approximate height for subtitle
    total_height = title_height + (subtitle_height if subtitle_text else 0)
    
    # Calculate top position to center the whole block
    left = Inches(0.6)
    width = Inches(9)
    top = (slide_height - total_height) / 2  # Center vertically on slide
    
    # Create a single text box for both title and subtitle
    text_box = slide.shapes.add_textbox(left, top, width, total_height)
    text_frame = text_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add title
    title = text_frame.add_paragraph()
    title.text = title_text
    title.font.name = TITLE_FONT
    title.font.size = TITLE_FONT_SIZE
    title.font.color.rgb = WHITE
    title.alignment = PP_ALIGN.LEFT
    
    # Add subtitle if provided
    if subtitle_text:
        subtitle = text_frame.add_paragraph()
        subtitle.text = subtitle_text
        subtitle.font.name = BODY_FONT
        subtitle.font.size = SUBTITLE_FONT_SIZE
        subtitle.font.color.rgb = WHITE
        subtitle.alignment = PP_ALIGN.LEFT
        # Add some space after the title
        title.space_after = Pt(12)

def apply_jep_slide_layout(slide, jep_number, jep_title, jep_description=None, slide_height=Inches(5.625)):
    """Apply layout for JEP slide with centered JEP number and title in a single text box"""
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ORANGE
    
    # Clear any existing placeholders first
    for shape in slide.shapes:
        if not (shape.has_text_frame and shape.text.strip()):
            sp = shape._element
            sp.getparent().remove(sp)
    
    # Calculate total height needed for JEP number and title
    jep_number_height = Inches(1.0)  # Approximate height for JEP number
    title_height = Inches(1.5)       # Approximate height for title
    spacing = Inches(0.3)            # Space between number and title
    total_height = jep_number_height + spacing + title_height
    
    # Calculate top position to center the whole block
    left = Inches(0.6)
    width = Inches(9)
    top = (slide_height - total_height) / 2
    
    # Create a single text box for both JEP number and title
    text_box = slide.shapes.add_textbox(left, top, width, total_height)
    text_frame = text_box.text_frame
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Add JEP number
    number = text_frame.add_paragraph()
    number.text = f"JEP {jep_number}"
    number.font.name = TITLE_FONT
    number.font.size = TITLE_FONT_SIZE
    number.font.color.rgb = WHITE
    number.alignment = PP_ALIGN.LEFT
    
    # Add JEP title
    title = text_frame.add_paragraph()
    title.text = jep_title
    title.font.name = BODY_FONT
    title.font.size = SUBTITLE_FONT_SIZE
    title.font.color.rgb = WHITE
    title.alignment = PP_ALIGN.LEFT
    number.space_after = Pt(12)

def apply_example_slide_layout(slide, title, content, slide_height=Inches(5.625)):
    """Apply layout for example slide with auto-fitting code"""
    # Set background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = ORANGE
    
    # Clear any existing placeholders first
    for shape in slide.shapes:
        if not (shape.has_text_frame and shape.text.strip()):
            sp = shape._element
            sp.getparent().remove(sp)
    
    # Add title at the top
    title_shape = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(9), Inches(1))
    title_shape.text = title
    title_frame = title_shape.text_frame
    
    # Style title
    title_frame.paragraphs[0].font.name = TITLE_FONT
    title_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_frame.paragraphs[0].font.color.rgb = WHITE
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    # Add content in a separate text box below the title
    content_box = slide.shapes.add_textbox(
        left=Inches(0.6),
        top=Inches(1.5),
        width=Inches(9),
        height=Inches(3.5)
    )
    
    # Set text frame properties
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    # Add content as a single paragraph with a monospace font for code
    p = text_frame.add_paragraph()
    p.text = content
    p.font.name = 'Courier New'  # Monospace font for code
    p.font.size = Pt(16)  # Fixed size that should work for most cases
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    p.font.bold = True
